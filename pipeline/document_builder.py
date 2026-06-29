"""
ProSEIT Innovation Agent — Document Builder

Converts Claude-generated markdown stage outputs into branded:
  • Word documents (.docx)     — all stages
  • Excel workbooks (.xlsx)    — data/table-heavy stages
  • PowerPoint slides (.pptx)  — key presentation stages
  • UIX prototype ZIP          — Stage 06 UIX Prototype
  • ZIP archive                — bundling all formats
"""
from __future__ import annotations
import io
import re
import zipfile
from datetime import date as _date

# ── openpyxl: top-level import (needed for module-level _XL_HDR constant) ────
from openpyxl import Workbook
from openpyxl.styles import (
    Font as XLFont, PatternFill as XLFill,
    Alignment as XLAlign, Border as XLBorder, Side as XLSide,
)
from openpyxl.utils import get_column_letter

# ── Lazy imports: docx and pptx deferred to first use ────────────────────────
def _docx_imports():
    global Document, Inches, Pt, WRGBColor, Cm, WD_ALIGN_PARAGRAPH, qn, OxmlElement
    global NAVY, _WN, _WG, _WW, _WGR, _WDT
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor as WRGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    NAVY = WRGBColor(0x24, 0x36, 0x4B)
    _WN = NAVY
    _WG = WRGBColor(0xD4, 0xA6, 0x2A)
    _WW = WRGBColor(0xFF, 0xFF, 0xFF)
    _WGR = WRGBColor(0xF4, 0xF6, 0xF8)
    _WDT = WRGBColor(0x1F, 0x29, 0x33)

def _pptx_imports():
    global Presentation, PIn, PPt, PRGBColor, PP_ALIGN
    global _PN, _PG, _PW, _PDT, _PGR
    from pptx import Presentation
    from pptx.util import Inches as PIn, Pt as PPt
    from pptx.dml.color import RGBColor as PRGBColor
    from pptx.enum.text import PP_ALIGN
    _PN = PRGBColor(0x24, 0x36, 0x4B)
    _PG = PRGBColor(0xD4, 0xA6, 0x2A)
    _PW = PRGBColor(0xFF, 0xFF, 0xFF)
    _PDT = PRGBColor(0x1F, 0x29, 0x33)
    _PGR = PRGBColor(0x88, 0x88, 0x88)

def _cairo_imports():
    global _cairosvg, _HAS_CAIRO
    try:
        import cairosvg as _cairosvg
        _HAS_CAIRO = True
    except Exception:
        _HAS_CAIRO = False

# Sentinels for lazy-loaded names
Document = Inches = Pt = WRGBColor = Cm = WD_ALIGN_PARAGRAPH = qn = OxmlElement = None
NAVY = _WN = _WG = _WW = _WGR = _WDT = None
Presentation = PIn = PPt = PRGBColor = PP_ALIGN = None
_PN = _PG = _PW = _PDT = _PGR = None
_cairosvg = None
_HAS_CAIRO = False

PROSEIT_CONTACT = (
    "ProSEIT · Plot 5 Blue Heights Plaza, Nkrumah Road, Kampala, Uganda"
    "  |  Tel: +256 790 334 653  |  www.proseit.org  |  info@proseit.org"
)
_TODAY = _date.today().strftime("%d %B %Y")

# ── ProSEIT logo SVG (background rect removed — transparent bg) ──────────────────
# White elements visible on navy cover; gold accent always visible.
_LOGO_SVG = b"""<?xml version=\"1.0\" encoding=\"UTF-8\"?>
<svg id=\"Layer_1\" data-name=\"Layer 1\" xmlns=\"http://www.w3.org/2000/svg\" viewBox=\"0 0 1080 1080\">
  <defs>
    <style>
      .cls-1 { fill: #fff; }
      .cls-3 { fill: #c8982c; }
    </style>
  </defs>
  <g>
    <g>
      <g>
        <path class=\"cls-3\" d=\"M490.01,457.8h18.71c.62,0,1.17-.39,1.38-.97l35.83-99.06,37.49-103.62c.3-.83-.31-1.7-1.2-1.7h-18.85c-.54,0-1.01.34-1.2.84l-35.31,97.59-38.07,105.22c-.3.83.31,1.7,1.2,1.7Z\"/>
        <path class=\"cls-1\" d=\"M622.44,206.25h-190.49c-1.1,0-1.99.89-1.99,1.99v270.1c0,.86.55,1.62,1.36,1.89l104.38,34.79c.41.14.85.14,1.26,0l104.38-34.79c.81-.27,1.36-1.03,1.36-1.89V208.25c0-1.1-.89-1.99-1.99-1.99h-18.27ZM622.44,297.43v167.41c0,.86-.55,1.62-1.36,1.89l-84.12,28.04c-.41.14-.85.14-1.26,0l-84.12-28.04c-.81-.27-1.36-1.03-1.36-1.89v-236.33c0-1.1.89-1.99,1.99-1.99h168.24c1.1,0,1.99.89,1.99,1.99v68.92Z\"/>
      </g>
      <g>
        <path class=\"cls-1\" d=\"M293.09,572.52c23.79,0,38.32,14.06,38.32,34.77s-14.52,34.92-38.32,34.92h-21.39c-7.08,0-5.81-1.34-5.81,5.81v27.48c0,.3-.24.55-.55.55h-19.92c-.3,0-.55-.24-.55-.55v-102.43c0-.3.24-.55.55-.55h47.66ZM293.09,623.2c11.28,0,17.46-6.33,17.46-15.91s-6.18-15.76-17.46-15.76h-26.65c-.3,0-.55.24-.55.55v30.58c0,.3.24.55.55.55h26.65Z\"/>
        <path class=\"cls-1\" d=\"M401.06,598.91l-1.03,18.68c-.02.29-.26.52-.55.52h-3.96c-16.7.23-27.18,6.18-27.5,33.07v24.33c0,.3-.24.55-.55.55h-19.3c-.3,0-.55-.24-.55-.55v-76c0-.3.24-.55.55-.55h19.3c.3,0,.55.24.55.55h0c0,8.75,1.52,9.6,8.68,4.4,5.16-3.75,11.41-6.03,18.98-6.03,1.86,0,3.32,0,4.99.49.24.07.4.31.39.56Z\"/>
        <path class=\"cls-1\" d=\"M490.55,637.57c0,24.26-19.01,40.48-41.56,40.48s-42.03-16.23-42.03-40.48,19-40.48,42.03-40.48,41.56,16.22,41.56,40.48ZM470.78,637.57c0-14.37-11.13-23.02-21.79-23.02s-22.25,8.65-22.25,23.02,11.12,23.02,22.25,23.02,21.79-8.65,21.79-23.02Z\"/>
        <path class=\"cls-1\" d=\"M582.79,600.14l-19.11.08c-.3.06-.58-.14-.64-.44-1.77-9.42-9.34-11.5-19.49-11.5s-16.84,5.87-16.84,12.82c0,5.25,2.32,8.96,10.35,10.51l22.25,4.48c16.84,3.56,27.19,13.6,27.19,29.67,0,20.09-16.53,32.6-41.1,32.6-21.45,0-43.2-12.34-43.19-31.84,0-.29.15-.56.44-.62l21.88-1.06c.29-.06.56.12.63.4,2.91,11.38,10.68,14.73,22.25,14.73s18.39-5.72,18.39-13.13c0-5.87-2.78-9.43-12.05-11.28l-21.94-4.17c-12.36-2.63-27.19-9.43-27.19-29.67,0-18.54,15.14-31.52,39.4-31.52,20.81,0,37.78,9.92,39.21,29.27.02.31-.13.6-.43.66Z\"/>
        <path class=\"cls-1\" d=\"M681.58,676.04h-73.19c-1.85,0-3.35-1.5-3.35-3.35v-96.17c0-2.21,1.79-4,4-4h72.23c1.16,0,2.11.94,2.11,2.11v14.79c0,1.16-.94,2.11-2.11,2.11h-54.67c-.3,0-.55.24-.55.55v22.24c0,.3.24.55.55.55h50.19c1.16,0,2.11.94,2.11,2.11v14.48c0,1.16-.94,2.11-2.11,2.11h-50.19c-.3,0-.55.24-.55.55v22.39c0,.3.24.55.55.55h54.98c1.16,0,2.11.94,2.11,2.11v14.79c0,1.16-.94,2.11-2.11,2.11Z\"/>
        <path class=\"cls-1\" d=\"M703,673.77v-98.97c0-1.26,1.02-2.28,2.28-2.28h16.46c1.26,0,2.28,1.02,2.28,2.28v98.97c0,1.26-1.02,2.28-2.28,2.28h-16.46c-1.26,0-2.28-1.02-2.28-2.28Z\"/>
        <path class=\"cls-1\" d=\"M793.38,592.84v81.9c0,.72-.59,1.31-1.31,1.31h-18.4c-.72,0-1.31-.59-1.31-1.31v-81.9c0-.72-.59-1.31-1.31-1.31h-31.84c-.72,0-1.31-.59-1.31-1.31v-16.39c0-.72.59-1.31,1.31-1.31h87.3c.72,0,1.31.59,1.31,1.31v16.39c0,.72-.59,1.31-1.31,1.31h-31.84c-.72,0-1.31.59-1.31,1.31Z\"/>
      </g>
    </g>
    <g>
      <g>
        <path class=\"cls-1\" d=\"M254.69,739.26v24.66h-6.13v-24.66h-10.05v-5.55h26.24v5.55h-10.05Z\"/>
        <path class=\"cls-1\" d=\"M290.84,750.22v13.7h-5.95v-12.44c0-3.16-1.8-5.27-4.55-5.27-3.38,0-5.59,2.25-5.59,7.75v9.96h-5.95v-32.46h5.95v11.99c1.76-1.76,4.1-2.71,7.08-2.71,5.41,0,9.02,3.88,9.02,9.47Z\"/>
        <path class=\"cls-1\" d=\"M317.84,754.27h-17.13c.59,3.25,2.89,5.23,5.64,5.23,1.76,0,3.97-.23,5.59-2.98l5.32,1.13c-1.98,4.69-5.99,6.94-10.91,6.94-6.36,0-11.54-4.82-11.54-11.9s5.18-11.95,11.63-11.95c5.99,0,11.18,4.64,11.4,11.49v2.03ZM300.89,750.08h10.86c-.77-2.89-2.93-4.15-5.32-4.15-2.25,0-4.82,1.35-5.55,4.15Z\"/>
        <path class=\"cls-1\" d=\"M363.19,746.29c.09,1.04.14,1.89.14,2.61,0,8.97-6.4,15.69-15.33,15.69s-16.05-6.9-16.05-15.78,6.99-15.78,15.87-15.78c6.36,0,12.35,2.98,14.52,9.11l-5.5,1.94c-2.21-4.24-5.72-5.27-9.02-5.27-5.41,0-9.65,4.33-9.65,10.01s4.37,10.05,10.01,10.05c4.6,0,8.02-2.7,9.06-6.99h-10.59v-5.59h16.54Z\"/>
        <path class=\"cls-1\" d=\"M389.74,763.92h-5.95v-2.03c-1.76,1.76-4.1,2.7-7.08,2.7-5.41,0-9.02-3.88-9.02-9.47v-13.7h5.95v12.44c0,3.16,1.8,5.27,4.55,5.27,3.38,0,5.59-2.25,5.59-7.75v-9.96h5.95v22.49Z\"/>
        <path class=\"cls-1\" d=\"M402.95,734.71c0,1.98-1.53,3.47-3.61,3.47s-3.61-1.49-3.61-3.47,1.49-3.56,3.61-3.56,3.61,1.67,3.61,3.56ZM396.37,763.92v-22.49h5.95v22.49h-5.95Z\"/>
        <path class=\"cls-1\" d=\"M408.94,763.92v-32.46h5.95v32.46h-5.95Z\"/>
        <path class=\"cls-1\" d=\"M443.83,763.92h-5.95v-1.76c-1.89,1.53-4.37,2.43-7.44,2.43-5.72,0-11.04-4.82-11.04-11.9s5.32-11.95,11.04-11.95c3.06,0,5.54.9,7.44,2.48v-11.76h5.95v32.46ZM437.88,752.69c0-4.15-3.2-6.76-6.4-6.76-3.61,0-6.31,2.61-6.31,6.76s2.7,6.72,6.31,6.72c3.2,0,6.4-2.57,6.4-6.72Z\"/>
        <path class=\"cls-1\" d=\"M483.68,752.69c0,7.08-5.55,11.81-12.13,11.81s-12.26-4.73-12.26-11.81,5.54-11.81,12.26-11.81,12.13,4.73,12.13,11.81ZM477.91,752.69c0-4.19-3.25-6.72-6.36-6.72s-6.49,2.52-6.49,6.72,3.25,6.72,6.49,6.72,6.36-2.52,6.36-6.72Z\"/>
        <path class=\"cls-1\" d=\"M501.67,741.43v4.87h-6.81v17.63h-5.95v-17.63h-3.43v-4.87h3.43v-2.03c0-5.27,3.61-8.52,8.88-8.52,1.44,0,2.75.27,4.01.86l-1.35,4.82c-.81-.27-1.62-.41-2.25-.41-1.94,0-3.34,1.13-3.34,3.25v2.03h6.81Z\"/>
        <path class=\"cls-1\" d=\"M529.39,733.72c6.94,0,11.18,4.1,11.18,10.14s-4.24,10.19-11.18,10.19h-7.93v9.87h-6.13v-30.2h14.06ZM529.39,748.5c3.29,0,5.09-1.85,5.09-4.64s-1.8-4.6-5.09-4.6h-7.93v9.24h7.93Z\"/>
        <path class=\"cls-1\" d=\"M560.9,741.29l-.32,5.72h-1.31c-5.5,0-8.02,3.38-8.02,9.65v7.26h-5.95v-22.49h5.95v4.1c1.8-2.61,4.42-4.42,8.07-4.42.59,0,1.04,0,1.58.18Z\"/>
        <path class=\"cls-1\" d=\"M587,752.69c0,7.08-5.55,11.81-12.13,11.81s-12.26-4.73-12.26-11.81,5.54-11.81,12.26-11.81,12.13,4.73,12.13,11.81ZM581.23,752.69c0-4.19-3.25-6.72-6.36-6.72s-6.49,2.52-6.49,6.72,3.25,6.72,6.49,6.72,6.36-2.52,6.36-6.72Z\"/>
        <path class=\"cls-1\" d=\"M604.99,741.43v4.87h-6.81v17.63h-5.95v-17.63h-3.43v-4.87h3.43v-2.03c0-5.27,3.61-8.52,8.88-8.52,1.44,0,2.75.27,4.01.86l-1.35,4.82c-.81-.27-1.62-.41-2.25-.41-1.94,0-3.34,1.13-3.34,3.25v2.03h6.81Z\"/>
        <path class=\"cls-1\" d=\"M628.92,754.27h-17.13c.59,3.25,2.89,5.23,5.64,5.23,1.76,0,3.97-.23,5.59-2.98l5.32,1.13c-1.98,4.69-5.99,6.94-10.91,6.94-6.36,0-11.54-4.82-11.54-11.9s5.18-11.95,11.63-11.95c5.99,0,11.18,4.64,11.4,11.49v2.03ZM611.97,750.08h10.86c-.77-2.89-2.93-4.15-5.32-4.15-2.25,0-4.82,1.35-5.55,4.15Z\"/>
        <path class=\"cls-1\" d=\"M637.53,756.75c.32,2.3,2.3,3.07,4.64,3.07,2.21,0,3.52-.99,3.52-2.25,0-.81-.54-1.62-2.48-2.03l-4.6-.95c-4.06-.86-6.36-3.16-6.36-6.4,0-4.33,3.52-7.44,9.2-7.44,5.32,0,9.02,2.75,9.78,6.67l-5.68,1.13c-.23-1.8-1.8-3.11-4.28-3.11s-3.06,1.17-3.06,2.21c0,.72.27,1.49,2.03,1.89l5.36,1.17c3.97.86,5.9,3.34,5.9,6.36,0,4.96-4.01,7.53-9.87,7.53-4.96,0-9.51-1.94-10.1-6.63l5.99-1.22Z\"/>
        <path class=\"cls-1\" d=\"M659.76,756.75c.32,2.3,2.3,3.07,4.64,3.07,2.21,0,3.52-.99,3.52-2.25,0-.81-.54-1.62-2.48-2.03l-4.6-.95c-4.06-.86-6.36-3.16-6.36-6.4,0-4.33,3.52-7.44,9.2-7.44,5.32,0,9.02,2.75,9.78,6.67l-5.68,1.13c-.23-1.8-1.8-3.11-4.28-3.11s-3.06,1.17-3.06,2.21c0,.72.27,1.49,2.03,1.89l5.36,1.17c3.97.86,5.9,3.34,5.9,6.36,0,4.96-4.01,7.53-9.87,7.53-4.96,0-9.51-1.94-10.1-6.63l5.99-1.22Z\"/>
        <path class=\"cls-1\" d=\"M685.23,734.71c0,1.98-1.53,3.47-3.61,3.47s-3.61-1.49-3.61-3.47,1.49-3.56,3.61-3.56,3.61,1.67,3.61,3.56ZM678.65,763.92v-22.49h5.95v22.49h-5.95Z\"/>
        <path class=\"cls-1\" d=\"M713.81,752.69c0,7.08-5.55,11.81-12.13,11.81s-12.26-4.73-12.26-11.81,5.54-11.81,12.26-11.81,12.13,4.73,12.13,11.81ZM708.04,752.69c0-4.19-3.25-6.72-6.36-6.72s-6.49,2.52-6.49,6.72,3.25,6.72,6.49,6.72,6.36-2.52,6.36-6.72Z\"/>
        <path class=\"cls-1\" d=\"M740.36,750.22v13.7h-5.95v-12.44c0-3.16-1.8-5.27-4.55-5.27-3.38,0-5.59,2.25-5.59,7.75v9.96h-5.95v-22.49h5.95v2.03c1.76-1.76,4.1-2.71,7.08-2.71,5.41,0,9.02,3.88,9.02,9.47Z\"/>
        <path class=\"cls-1\" d=\"M768.71,763.92h-5.95v-1.76c-1.89,1.53-4.37,2.43-7.44,2.43-5.72,0-11.04-4.82-11.04-11.9s5.32-11.95,11.04-11.95c3.06,0,5.54.9,7.44,2.48v-1.8h5.95v22.49ZM762.76,752.69c0-4.15-3.2-6.76-6.4-6.76-3.61,0-6.31,2.61-6.31,6.76s2.7,6.72,6.31,6.72c3.2,0,6.4-2.57,6.4-6.72Z\"/>
        <path class=\"cls-1\" d=\"M775.02,763.92v-32.46h5.95v32.46h-5.95Z\"/>
        <path class=\"cls-1\" d=\"M263.79,795.83l-6.13,1.26c-.45-2.84-2.93-4.69-5.95-4.69s-4.91,1.71-4.91,3.74c0,1.53.68,2.61,3.02,3.07l6.49,1.31c4.91,1.04,7.93,3.97,7.93,8.66,0,5.86-4.82,9.51-11.99,9.51-6.31,0-11.86-2.93-12.62-8.97l6.45-1.31c.81,3.43,3.34,4.91,6.76,4.91s5.36-1.67,5.36-3.83c0-1.71-.81-2.75-3.52-3.29l-6.4-1.22c-3.61-.77-7.93-2.75-7.93-8.66,0-5.41,4.42-9.2,11.49-9.2,6.13,0,10.77,3.34,11.95,8.7Z\"/>
        <path class=\"cls-1\" d=\"M292.24,806.79c0,7.08-5.55,11.81-12.13,11.81s-12.26-4.73-12.26-11.81,5.54-11.81,12.26-11.81,12.13,4.73,12.13,11.81ZM286.47,806.79c0-4.19-3.25-6.72-6.36-6.72s-6.49,2.52-6.49,6.72,3.25,6.72,6.49,6.72,6.36-2.52,6.36-6.72Z\"/>
        <path class=\"cls-1\" d=\"M310.22,795.52v4.87h-6.81v17.63h-5.95v-17.63h-3.43v-4.87h3.43v-2.03c0-5.27,3.61-8.52,8.88-8.52,1.44,0,2.75.27,4.01.86l-1.35,4.82c-.81-.27-1.62-.41-2.25-.41-1.94,0-3.34,1.13-3.34,3.25v2.03h6.81Z\"/>
        <path class=\"cls-1\" d=\"M330.64,817.02c-2.34,1.26-3.97,1.67-5.86,1.67-5.32,0-8.84-2.79-8.84-8.97v-9.33h-4.82v-4.87h4.82v-6.81h5.95v6.81h7.35v4.87h-7.35v8.88c0,2.71,1.31,3.88,3.29,3.88,1.04,0,2.57-.45,3.74-1.08l1.71,4.96Z\"/>
        <path class=\"cls-1\" d=\"M358.23,818.06h-5.36l-4.42-13.25-4.46,13.25h-5.36l-7.48-22.54h5.91l4.46,13.84,4.6-13.84h4.64l4.6,13.88,4.46-13.88h5.91l-7.48,22.54Z\"/>
        <path class=\"cls-1\" d=\"M391.72,818.01h-5.95v-1.76c-1.89,1.53-4.37,2.43-7.44,2.43-5.72,0-11.04-4.82-11.04-11.9s5.32-11.95,11.04-11.95c3.06,0,5.54.9,7.44,2.48v-1.8h5.95v22.49ZM385.77,806.79c0-4.15-3.2-6.76-6.4-6.76-3.61,0-6.31,2.61-6.31,6.76s2.7,6.72,6.31,6.72c3.2,0,6.4-2.57,6.4-6.72Z\"/>
        <path class=\"cls-1\" d=\"M413.63,795.38l-.32,5.72h-1.31c-5.5,0-8.02,3.38-8.02,9.65v7.26h-5.95v-22.49h5.95v4.1c1.8-2.61,4.42-4.42,8.07-4.42.59,0,1.04,0,1.58.18Z\"/>
        <path class=\"cls-1\" d=\"M438.38,808.37h-17.13c.59,3.25,2.89,5.23,5.64,5.23,1.76,0,3.97-.23,5.59-2.98l5.32,1.13c-1.98,4.69-5.99,6.94-10.91,6.94-6.36,0-11.54-4.82-11.54-11.9s5.18-11.95,11.63-11.95c5.99,0,11.18,4.64,11.4,11.49v2.03ZM421.43,804.17h10.86c-.77-2.89-2.93-4.15-5.32-4.15-2.25,0-4.82,1.35-5.55,4.15Z\"/>
        <path class=\"cls-1\" d=\"M476.79,818.01h-22.95v-30.2h22.86v5.55h-16.72v6.81h15.42v5.45h-15.42v6.85h16.81v5.55Z\"/>
        <path class=\"cls-1\" d=\"M504.46,804.31v13.7h-5.95v-12.44c0-3.16-1.8-5.27-4.55-5.27-3.38,0-5.59,2.25-5.59,7.75v9.96h-5.95v-22.49h5.95v2.03c1.76-1.76,4.1-2.71,7.08-2.71,5.41,0,9.02,3.88,9.02,9.47Z\"/>
        <path class=\"cls-1\" d=\"M532.14,815.85c0,6.9-4.78,10.73-12.04,10.73-5.23,0-9.02-1.85-10.95-4.96l4.82-3.65c1.76,2.21,3.2,3.25,6.13,3.25,3.79,0,6.27-2.34,6.27-6.27v-.36c-1.67,1.35-3.97,2.07-7.08,2.07-5.63,0-10.91-4.78-10.91-10.95s5.27-10.86,10.91-10.86c3.11,0,5.41.77,7.08,2.07v-1.4h5.72l.05,20.33ZM526.37,805.71c0-3.52-2.89-5.86-6.04-5.86-3.56,0-6.22,2.34-6.22,5.86s2.66,5.86,6.22,5.86c3.16,0,6.04-2.34,6.04-5.86Z\"/>
        <path class=\"cls-1\" d=\"M545.35,788.8c0,1.98-1.53,3.47-3.61,3.47s-3.61-1.49-3.61-3.47,1.49-3.56,3.61-3.56,3.61,1.67,3.61,3.56ZM538.77,818.01v-22.49h5.95v22.49h-5.95Z\"/>
        <path class=\"cls-1\" d=\"M573.39,804.31v13.7h-5.95v-12.44c0-3.16-1.8-5.27-4.55-5.27-3.38,0-5.59,2.25-5.59,7.75v9.96h-5.95v-22.49h5.95v2.03c1.76-1.76,4.1-2.71,7.08-2.71,5.41,0,9.02,3.88,9.02,9.47Z\"/>
        <path class=\"cls-1\" d=\"M600.35,808.37h-17.13c.59,3.25,2.89,5.23,5.64,5.23,1.76,0,3.97-.23,5.59-2.98l5.32,1.13c-1.98,4.69-5.99,6.94-10.91,6.94-6.36,0-11.54-4.82-11.54-11.9s5.18-11.95,11.63-11.95c5.99,0,11.18,4.64,11.4,11.49v2.03ZM583.4,804.17h10.86c-.77-2.89-2.93-4.15-5.32-4.15-2.25,0-4.82,1.35-5.55,4.15Z\"/>
        <path class=\"cls-1\" d=\"M626.09,808.37h-17.13c.59,3.25,2.89,5.23,5.64,5.23,1.76,0,3.97-.23,5.59-2.98l5.32,1.13c-1.98,4.69-5.99,6.94-10.91,6.94-6.36,0-11.54-4.82-11.54-11.9s5.18-11.95,11.63-11.95c5.99,0,11.18,4.64,11.4,11.49v2.03ZM609.14,804.17h10.86c-.77-2.89-2.93-4.15-5.32-4.15-2.25,0-4.82,1.35-5.55,4.15Z\"/>
        <path class=\"cls-1\" d=\"M646.19,795.38l-.32,5.72h-1.31c-5.5,0-8.02,3.38-8.02,9.65v7.26h-5.95v-22.49h5.95v4.1c1.8-2.61,4.42-4.42,8.07-4.42.59,0,1.04,0,1.58.18Z\"/>
        <path class=\"cls-1\" d=\"M653.81,810.85c.32,2.3,2.3,3.07,4.64,3.07,2.21,0,3.52-.99,3.52-2.25,0-.81-.54-1.62-2.48-2.03l-4.6-.95c-4.06-.86-6.36-3.16-6.36-6.4,0-4.33,3.52-7.44,9.2-7.44,5.32,0,9.02,2.75,9.78,6.67l-5.68,1.13c-.23-1.8-1.8-3.11-4.28-3.11s-3.06,1.17-3.06,2.21c0,.72.27,1.49,2.03,1.89l5.36,1.17c3.97.86,5.9,3.34,5.9,6.36,0,4.96-4.01,7.53-9.87,7.53-4.96,0-9.51-1.94-10.1-6.63l5.99-1.22Z\"/>
        <path class=\"cls-1\" d=\"M703.26,815.76c-2.39,1.98-5.5,2.93-8.79,2.93-7.21,0-12.04-3.79-12.04-9.56,0-3.34,2.03-5.95,5.14-7.62-1.49-1.85-2.16-3.74-2.16-6.04,0-4.96,4.15-8.34,9.6-8.34,5.72,0,8.97,3.56,9.65,8.61l-5.59.99c-.22-3.38-1.76-4.6-4.06-4.6-2.16,0-3.88,1.13-3.88,3.47,0,1.89.9,3.25,3.92,5.86l7.17,6.27c.27-.95.41-1.98.41-3.16v-2.3h5.27v2.3c0,2.7-.5,5-1.4,6.94l4.33,3.79-3.52,4.01-4.06-3.56Z\"/>
        <path class=\"cls-1\" d=\"M240.76,872.11v-30.2h6.13v30.2h-6.13Z\"/>
        <path class=\"cls-1\" d=\"M267.13,847.45v24.66h-6.13v-24.66h-10.05v-5.55h26.24v5.55h-10.05Z\"/>
        <path class=\"cls-1\" d=\"M306.25,841.9c6.94,0,11.18,4.1,11.18,10.14s-4.24,10.19-11.18,10.19h-7.93v9.87h-6.13v-30.2h14.06ZM306.25,856.69c3.29,0,5.09-1.85,5.09-4.64s-1.8-4.6-5.09-4.6h-7.93v9.24h7.93Z\"/>
        <path class=\"cls-1\" d=\"M337.76,849.48l-.32,5.72h-1.31c-5.5,0-8.02,3.38-8.02,9.65v7.26h-5.95v-22.49h5.95v4.1c1.8-2.61,4.42-4.42,8.07-4.42.59,0,1.04,0,1.58.18Z\"/>
        <path class=\"cls-1\" d=\"M363.91,872.11h-5.95v-1.76c-1.89,1.53-4.37,2.43-7.44,2.43-5.72,0-11.04-4.82-11.04-11.9s5.32-11.95,11.04-11.95c3.06,0,5.54.9,7.44,2.48v-1.8h5.95v22.49ZM357.96,860.88c0-4.15-3.2-6.76-6.4-6.76-3.61,0-6.31,2.61-6.31,6.76s2.7,6.72,6.31,6.72c3.2,0,6.4-2.57,6.4-6.72Z\"/>
        <path class=\"cls-1\" d=\"M386.22,862.96l5.5,1.4c-1.26,5.14-5.91,8.43-11.27,8.43-6.58,0-12.04-4.82-12.04-11.9s5.45-11.95,12.04-11.95c5.23,0,9.78,3.25,11.27,8.29l-5.72,1.62c-.81-3.16-2.93-4.46-5.54-4.46-3.74,0-6.27,2.66-6.27,6.49s2.52,6.45,6.27,6.45c2.61,0,4.73-1.26,5.77-4.37Z\"/>
        <path class=\"cls-1\" d=\"M412.59,871.12c-2.34,1.26-3.97,1.67-5.86,1.67-5.32,0-8.84-2.79-8.84-8.97v-9.33h-4.82v-4.87h4.82v-6.81h5.95v6.81h7.35v4.87h-7.35v8.88c0,2.71,1.31,3.88,3.29,3.88,1.04,0,2.57-.45,3.74-1.08l1.71,4.96Z\"/>
        <path class=\"cls-1\" d=\"M422.69,842.9c0,1.98-1.53,3.47-3.61,3.47s-3.61-1.49-3.61-3.47,1.49-3.56,3.61-3.56,3.61,1.67,3.61,3.56ZM416.11,872.11v-22.49h5.95v22.49h-5.95Z\"/>
        <path class=\"cls-1\" d=\"M445.05,871.12c-2.34,1.26-3.97,1.67-5.86,1.67-5.32,0-8.84-2.79-8.84-8.97v-9.33h-4.82v-4.87h4.82v-6.81h5.95v6.81h7.35v4.87h-7.35v8.88c0,2.71,1.31,3.88,3.29,3.88,1.04,0,2.57-.45,3.74-1.08l1.71,4.96Z\"/>
        <path class=\"cls-1\" d=\"M455.15,842.9c0,1.98-1.53,3.47-3.61,3.47s-3.61-1.49-3.61-3.47,1.49-3.56,3.61-3.56,3.61,1.67,3.61,3.56ZM448.57,872.11v-22.49h5.95v22.49h-5.95Z\"/>
        <path class=\"cls-1\" d=\"M483.73,860.88c0,7.08-5.55,11.81-12.13,11.81s-12.26-4.73-12.26-11.81,5.54-11.81,12.26-11.81,12.13,4.73,12.13,11.81ZM477.96,860.88c0-4.19-3.25-6.72-6.36-6.72s-6.49,2.52-6.49,6.72,3.25,6.72,6.49,6.72,6.36-2.52,6.36-6.72Z\"/>
        <path class=\"cls-1\" d=\"M510.28,858.4v13.7h-5.95v-12.44c0-3.16-1.8-5.27-4.55-5.27-3.38,0-5.59,2.25-5.59,7.75v9.96h-5.95v-22.49h5.95v2.03c1.76-1.76,4.1-2.71,7.08-2.71,5.41,0,9.02,3.88,9.02,9.47Z\"/>
        <path class=\"cls-1\" d=\"M537.23,862.46h-17.13c.59,3.25,2.89,5.23,5.64,5.23,1.76,0,3.97-.23,5.59-2.98l5.32,1.13c-1.98,4.69-5.99,6.94-10.91,6.94-6.36,0-11.54-4.82-11.54-11.9s5.18-11.95,11.63-11.95c6,0,11.18,4.64,11.41,11.49v2.03ZM520.29,858.27h10.86c-.77-2.89-2.93-4.15-5.32-4.15-2.25,0-4.82,1.35-5.54,4.15Z\"/>
        <path class=\"cls-1\" d=\"M557.34,849.48l-.32,5.72h-1.31c-5.5,0-8.02,3.38-8.02,9.65v7.26h-5.95v-22.49h5.95v4.1c1.8-2.61,4.42-4.42,8.07-4.42.59,0,1.04,0,1.58.18Z\"/>
        <path class=\"cls-1\" d=\"M564.96,864.94c.32,2.3,2.3,3.07,4.64,3.07,2.21,0,3.52-.99,3.52-2.25,0-.81-.54-1.62-2.48-2.03l-4.6-.95c-4.06-.86-6.36-3.16-6.36-6.4,0-4.33,3.52-7.44,9.2-7.44,5.32,0,9.02,2.75,9.78,6.67l-5.68,1.13c-.23-1.8-1.8-3.11-4.28-3.11s-3.06,1.17-3.06,2.21c0,.72.27,1.49,2.03,1.89l5.36,1.17c3.97.86,5.9,3.34,5.9,6.36,0,4.96-4.01,7.53-9.87,7.53-4.96,0-9.51-1.94-10.1-6.63l5.99-1.22Z\"/>
      </g>
      <path class=\"cls-3\" d=\"M773.25,873.75h14.71c.19,0,.36-.12.43-.3l25.87-71.53,27.21-75.21c.07-.21-.08-.42-.3-.42h-14.81c-.13,0-.25.08-.3.21l-25.5,70.48-27.63,76.36c-.07.21.08.42.3.42Z\"/>
    </g>
  </g>
</svg>"""


def _logo_png_bytes(width_px: int = 200) -> bytes | None:
    """Render logo SVG to PNG bytes. Returns None if cairosvg unavailable."""
    _cairo_imports()
    if not _HAS_CAIRO:
        return None
    try:
        return _cairosvg.svg2png(bytestring=_LOGO_SVG, output_width=width_px)
    except Exception:
        return None


# ── Stage routing ────────────────────────────────────────────────────────────────────────────────────
EXCEL_STAGES = frozenset({
    "03_editions", "05_user_stories", "06_ux_wireframes", "07_data_model",
    "08_api_design", "09_integrations", "11_testing",
    "14_budget", "17_roadmap", "18_package",
})
PPTX_STAGES = frozenset({
    "00_intake", "02_concept", "15_business_plan",
    "17_roadmap", "18_package",
})
UIX_STAGES = frozenset({"06_ux_wireframes"})

STAGE_DELIVERABLES: dict[str, list[str]] = {
    "00_intake":        ["Word (.docx)", "PowerPoint (.pptx)"],
    "01_policy":        ["Word (.docx)"],
    "02_concept":       ["Word (.docx)", "PowerPoint (.pptx)"],
    "03_editions":      ["Word (.docx)", "Excel (.xlsx)"],
    "04_architecture":  ["Word (.docx)"],
    "05_user_stories":  ["Word (.docx)", "Excel (.xlsx)"],
    "06_ux_wireframes": ["UIX Prototype Pack (.zip)", "Word Reference (.docx)", "Excel Mapping (.xlsx)"],
    "07_data_model":    ["Word (.docx)", "Excel (.xlsx)"],
    "08_api_design":    ["Word (.docx)", "Excel (.xlsx)"],
    "09_integrations":  ["Word (.docx)", "Excel (.xlsx)"],
    "10_security":      ["Word (.docx)"],
    "11_testing":       ["Word (.docx)", "Excel (.xlsx)"],
    "12_deployment":    ["Word (.docx)"],
    "13_documentation": ["Word (.docx)"],
    "14_budget":        ["Word (.docx)", "Excel Budget (.xlsx)"],
    "15_business_plan": ["Word (.docx)", "PowerPoint (.pptx)"],
    "16_governance":    ["Word (.docx)"],
    "17_roadmap":       ["Word (.docx)", "Excel (.xlsx)", "PowerPoint (.pptx)"],
    "18_package":       ["Word (.docx)", "PowerPoint (.pptx)", "Excel (.xlsx)"],
}

# ═══════════════════════════════════════════════════════════════════════════
# UIX Prototype file parser
# ═══════════════════════════════════════════════════════════════════════════

_UIX_FILE_RE = re.compile(
    r"===UIX_FILE:\s*([^\n]+?)\s*===\n(.*?)===UIX_END===",
    re.DOTALL,
)


def parse_uix_files(content: str) -> dict[str, str]:
    """Extract filename->content pairs from UIX delimiter blocks."""
    return {m.group(1).strip(): m.group(2) for m in _UIX_FILE_RE.finditer(content)}


def build_uix_zip(stage: str, stage_label: str, content: str, innov_title: str) -> bytes:
    """Build a browsable UIX prototype ZIP from Claude's delimited output."""
    files = parse_uix_files(content)
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        if files:
            for filename, file_content in files.items():
                zf.writestr(f"uix_prototype/{filename}", file_content)
        else:
            zf.writestr("uix_prototype/readme.md", content)

        launch_html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta http-equiv="refresh" content="0; url=i.html">
<title>{innov_title} — UIX Prototype</title>
<style>
  body {{ font-family: Calibri, sans-serif; background: #24364B; color: #fff;
         display: flex; align-items: center; justify-content: center;
         min-height: 100vh; margin: 0; }}
  .card {{ background: rgba(255,255,255,0.08); border: 2px solid #D4A62A;
           border-radius: 8px; padding: 40px 60px; text-align: center; }}
  h1 {{ color: #D4A62A; margin: 0 0 8px; }}
  p {{ opacity: 0.7; margin: 0 0 24px; }}
  a {{ display: inline-block; background: #D4A62A; color: #24364B;
       font-weight: 700; padding: 12px 28px; border-radius: 4px;
       text-decoration: none; }}
</style>
</head>
<body>
<div class="card">
  <h1>ProSEIT UIX Prototype</h1>
  <p>{innov_title}</p>
  <p style="font-size:0.85rem">{stage_label}</p>
  <a href="i.html">Open Prototype &rarr;</a>
</div>
</body>
</html>"""
        zf.writestr("uix_prototype/launch.html", launch_html)

        file_list = "\n".join(f"- {fn}" for fn in sorted(files.keys())) if files else "(no structured files found)"
        manifest = f"""# {innov_title} — UIX Prototype Pack

Generated: {_TODAY}
Stage: {stage_label}
{PROSEIT_CONTACT}

## Files
{file_list}

## How to use
1. Extract this ZIP to a folder
2. Open `launch.html` in a browser (or open `i.html` directly)
3. Navigate using the sidebar

## Notes
- All pages link to shared `a/s.css` and `a/j.js`
- No build step required — pure HTML/CSS/JS
- For Angular integration see `readme.md` inside the prototype folder
"""
        zf.writestr("MANIFEST.md", manifest)

    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════
# Shared markdown parser
# ═══════════════════════════════════════════════════════════════════════════

def _parse_md(text: str) -> list[tuple[str, str]]:
    items: list[tuple[str, str]] = []
    for raw in text.split("\n"):
        line = raw.rstrip()
        if not line:
            items.append(("blank", ""))
        elif line.startswith("#### "):
            items.append(("h4", line[5:]))
        elif line.startswith("### "):
            items.append(("h3", line[4:]))
        elif line.startswith("## "):
            items.append(("h2", line[3:]))
        elif line.startswith("# "):
            items.append(("h1", line[2:]))
        elif line.startswith("|") and "|" in line[1:]:
            items.append(("table_row", line))
        elif re.match(r"^\s*[-*]\s", line):
            items.append(("bullet", re.sub(r"^\s*[-*]\s", "", line)))
        elif line.startswith("> "):
            items.append(("quote", line[2:]))
        elif re.match(r"^[-=]{3,}$", line.strip()):
            items.append(("hr", ""))
        else:
            items.append(("para", line))
    return items


def _plain(text: str) -> str:
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"\*(.*?)\*", r"\1", text)
    text = re.sub(r"`(.*?)`", r"\1", text)
    text = re.sub(r"\[(.*?)\]\(.*?\)", r"\1", text)
    return text.strip()


def _extract_tables(text: str) -> list[list[list[str]]]:
    tables: list[list[list[str]]] = []
    current: list[list[str]] = []
    for t, content in _parse_md(text):
        if t == "table_row":
            cells = [c.strip() for c in content.strip("|").split("|")]
            if all(re.match(r"^[-: ]+$", c) for c in cells if c):
                continue
            current.append(cells)
        else:
            if current:
                tables.append(current)
                current = []
    if current:
        tables.append(current)
    return tables


# ═══════════════════════════════════════════════════════════════════════════
# Word document builder
# ═══════════════════════════════════════════════════════════════════════════

def _cell_bg(cell, hex6: str) -> None:
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex6)
    tcPr.append(shd)


def _add_cover(doc: Document, innov_title: str, stage_label: str) -> None:
    tbl = doc.add_table(rows=1, cols=1)
    tbl.style = "Table Grid"
    cell = tbl.cell(0, 0)
    _cell_bg(cell, "24364B")

    def _cp(text: str, size: int, color: WRGBColor,
            bold: bool = False, italic: bool = False) -> None:
        p = cell.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(text)
        r.bold = bold
        r.italic = italic
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"

    # Logo — top-left, small and sharp
    p_logo = cell.paragraphs[0]
    p_logo.alignment = WD_ALIGN_PARAGRAPH.LEFT
    logo_png = _logo_png_bytes(240)
    if logo_png:
        r_logo = p_logo.add_run()
        r_logo.add_picture(io.BytesIO(logo_png), width=Inches(1.4))
    else:
        # Fallback: text wordmark when cairosvg unavailable
        r_logo = p_logo.add_run("PROSEIT")
        r_logo.bold = True
        r_logo.font.size = Pt(14)
        r_logo.font.color.rgb = _WG
        r_logo.font.name = "Calibri"

    _cp(" ", 4, _WW)
    _cp("PROSEIT", 42, _WG, bold=True)
    _cp("Professional Society for Engineering, Innovation & Technology", 11, _WW)
    _cp("─" * 40, 10, _WG)
    _cp(" ", 6, _WW)
    _cp("INNOVATION PIPELINE REPORT", 12, _WG, bold=True)
    _cp(" ", 4, _WW)
    _cp(innov_title, 22, _WW, bold=True)
    _cp(" ", 4, _WW)
    _cp(stage_label, 14, _WG)
    _cp(" ", 6, _WW)
    _cp(_TODAY, 10, _WW)
    _cp(" ", 6, _WW)
    _cp("CONFIDENTIAL — ProSEIT Internal Document", 9, _WG, italic=True)
    _cp(" ", 6, _WW)

    tbl.rows[0].height = Cm(20)
    doc.add_page_break()


def _setup_hf(doc: Document, stage_label: str, innov_title: str) -> None:
    section = doc.sections[0]
    section.header_distance = Cm(1.0)
    section.footer_distance = Cm(0.8)

    hdr = section.header
    hdr.paragraphs[0].clear()
    p = hdr.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    r1 = p.add_run("ProSEIT  |  ")
    r1.bold = True
    r1.font.color.rgb = _WN
    r1.font.size = Pt(9)
    r2 = p.add_run(stage_label)
    r2.font.color.rgb = _WG
    r2.font.size = Pt(9)
    r3 = p.add_run(f"  —  {innov_title}")
    r3.font.color.rgb = _WN
    r3.font.size = Pt(8)

    ftr = section.footer
    ftr.paragraphs[0].clear()
    pf = ftr.paragraphs[0]
    pf.alignment = WD_ALIGN_PARAGRAPH.CENTER
    rf = pf.add_run(PROSEIT_CONTACT)
    rf.font.size = Pt(7.5)
    rf.font.color.rgb = _WN


def _add_runs(para, text: str, base_pt: int = 11) -> None:
    parts = re.split(r"(\*\*.*?\*\*|\*[^*]+\*|`[^`]+`)", text)
    for part in parts:
        if part.startswith("**") and part.endswith("**") and len(part) > 4:
            r = para.add_run(_plain(part))
            r.bold = True
        elif part.startswith("*") and part.endswith("*") and len(part) > 2:
            r = para.add_run(part[1:-1])
            r.italic = True
        elif part.startswith("`") and part.endswith("`") and len(part) > 2:
            r = para.add_run(part[1:-1])
            r.font.name = "Courier New"
            r.font.size = Pt(base_pt - 1)
        elif part:
            para.add_run(part)


def _render_to_docx(doc: Document, tokens: list[tuple[str, str]]) -> None:
    i = 0
    while i < len(tokens):
        t, content = tokens[i]

        if t == "h1":
            p = doc.add_heading(level=1)
            p.clear()
            r = p.add_run(_plain(content))
            r.font.color.rgb = _WN
            r.bold = True
        elif t == "h2":
            p = doc.add_heading(level=2)
            p.clear()
            r = p.add_run(_plain(content))
            r.font.color.rgb = _WN
        elif t == "h3":
            p = doc.add_heading(level=3)
            p.clear()
            r = p.add_run(_plain(content))
            r.font.color.rgb = _WG
            r.bold = True
        elif t == "h4":
            p = doc.add_heading(level=4)
            p.clear()
            r = p.add_run(_plain(content))
            r.font.color.rgb = _WDT
        elif t == "bullet":
            p = doc.add_paragraph(style="List Bullet")
            p.paragraph_format.left_indent = Cm(0.5)
            _add_runs(p, content)
        elif t == "quote":
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Cm(1.0)
            r = p.add_run(content)
            r.italic = True
            r.font.color.rgb = WRGBColor(0x55, 0x66, 0x77)
        elif t == "table_row":
            rows: list[list[str]] = []
            j = i
            while j < len(tokens) and tokens[j][0] == "table_row":
                cells = [c.strip() for c in tokens[j][1].strip("|").split("|")]
                if not all(re.match(r"^[-: ]+$", c) for c in cells if c):
                    rows.append(cells)
                j += 1
            i = j
            if rows:
                ncols = max(len(r) for r in rows)
                padded = [r + [""] * (ncols - len(r)) for r in rows]
                tbl = doc.add_table(rows=len(padded), cols=ncols)
                tbl.style = "Table Grid"
                for ri, row in enumerate(padded):
                    for ci, val in enumerate(row):
                        cell = tbl.cell(ri, ci)
                        cell.paragraphs[0].clear()
                        if ri == 0:
                            _cell_bg(cell, "24364B")
                            r = cell.paragraphs[0].add_run(_plain(val))
                            r.bold = True
                            r.font.color.rgb = _WW
                            r.font.size = Pt(9)
                        else:
                            if ri % 2 == 0:
                                _cell_bg(cell, "F4F6F8")
                            pr = cell.paragraphs[0].add_run(_plain(val))
                            pr.font.size = Pt(9)
                doc.add_paragraph()
            continue
        elif t == "hr":
            p = doc.add_paragraph("_" * 60)
            p.paragraph_format.space_after = Pt(4)
        elif t == "para" and content.strip():
            p = doc.add_paragraph()
            _add_runs(p, content)
            p.paragraph_format.space_after = Pt(4)

        i += 1


def build_word_doc(
    stage: str, stage_label: str, content: str, innov_title: str
) -> bytes:
    _docx_imports()
    doc = Document()
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)
    _add_cover(doc, innov_title, stage_label)
    _setup_hf(doc, stage_label, innov_title)
    _render_to_docx(doc, _parse_md(content))
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════
# Excel workbook builder
# ═══════════════════════════════════════════════════════════════════════════

_XL_HDR = {
    "font": XLFont(bold=True, color="FFFFFF", name="Calibri", size=10),
    "fill": XLFill(fill_type="solid", fgColor="24364B"),
    "alignment": XLAlign(horizontal="center", vertical="center", wrap_text=True),
    "border": XLBorder(
        left=XLSide(style="thin"), right=XLSide(style="thin"),
        top=XLSide(style="thin"), bottom=XLSide(style="thin"),
    ),
}


def _xl_row_style(ri: int) -> dict:
    bg = "F4F6F8" if ri % 2 == 0 else "FFFFFF"
    return {
        "font": XLFont(name="Calibri", size=10),
        "fill": XLFill(fill_type="solid", fgColor=bg),
        "alignment": XLAlign(horizontal="left", vertical="top", wrap_text=True),
        "border": XLBorder(
            left=XLSide(style="thin"), right=XLSide(style="thin"),
            top=XLSide(style="thin"), bottom=XLSide(style="thin"),
        ),
    }


def _xl_write_row(ws, row_num: int, cells: list[str], style: dict) -> None:
    for ci, val in enumerate(cells, 1):
        c = ws.cell(row=row_num, column=ci, value=_plain(val))
        for attr, v in style.items():
            setattr(c, attr, v)


def _build_uix_screen_map(content: str) -> list[list[str]]:
    """Extract file inventory from UIX content for the Excel mapping sheet."""
    files = parse_uix_files(content)
    screen_map = [
        ["File", "Screen / Purpose", "Angular Component", "Route", "DocType(s)"],
    ]
    screen_names = {
        "i.html": ("Dashboard / Home", "DashboardComponent", "/app/dashboard"),
        "p1.html": ("Feature Screen 1", "Feature1Component", "/app/feature/1"),
        "p2.html": ("Feature Screen 2", "Feature2Component", "/app/feature/2"),
        "p3.html": ("Feature Screen 3", "Feature3Component", "/app/feature/3"),
        "p4.html": ("Feature Screen 4", "Feature4Component", "/app/feature/4"),
        "p5.html": ("Feature Screen 5", "Feature5Component", "/app/feature/5"),
        "p6.html": ("List Screen", "ListComponent", "/app/list"),
        "p7.html": ("Detail / Form", "DetailComponent", "/app/detail/:id"),
        "p8.html": ("Reports / Analytics", "ReportsComponent", "/app/reports"),
        "g.html": ("Global Settings", "SettingsComponent", "/app/settings"),
        "s.html": ("Superadmin Panel", "AdminComponent", "/app/admin"),
        "a/s.css": ("Master Stylesheet", "N/A", "N/A"),
        "a/j.js": ("Master JavaScript", "N/A", "N/A"),
        "readme.md": ("Developer Handover", "N/A", "N/A"),
    }
    fname_list = sorted(files.keys()) if files else sorted(screen_names.keys())
    for fname in fname_list:
        info = screen_names.get(fname, (fname, "CustomComponent", "/app/custom"))
        screen_map.append([fname, info[0], info[1], info[2], "[see Developer Handover Panel]"])
    return screen_map


def build_excel_doc(
    stage: str, stage_label: str, content: str, innov_title: str
) -> bytes:
    wb = Workbook()
    wb.remove(wb.active)
    tables = _extract_tables(content)
    tokens = _parse_md(content)

    def _title_rows(ws) -> None:
        ws["A1"] = innov_title
        ws["A1"].font = XLFont(bold=True, size=13, color="24364B", name="Calibri")
        ws["A2"] = stage_label
        ws["A2"].font = XLFont(size=10, color="D4A62A", name="Calibri")
        ws["A3"] = PROSEIT_CONTACT
        ws["A3"].font = XLFont(size=8, color="888888", name="Calibri")
        ws.row_dimensions[1].height = 18

    if stage in UIX_STAGES:
        uix_files = parse_uix_files(content)
        xl_map_content = uix_files.get("xl_map.md", "")
        if xl_map_content:
            # Parse 13 sections from xl_map.md and create 13 Excel sheets
            sheet_order = [
                "UI Catalog", "Story Coverage", "Epic Coverage", "Angular Components",
                "Backend Mapping", "ERPNext Reuse", "Frappe DocTypes APIs", "QA Acceptance",
                "Reports", "Mobile Screens", "Image Use", "Project Structure", "Intern Resources",
            ]
            # Split by ## headings
            parts = re.split(r'^## (.+)', xl_map_content, flags=re.MULTILINE)
            section_data: dict[str, list[list[str]]] = {}
            for pi in range(1, len(parts), 2):
                heading = parts[pi].strip()
                body = parts[pi + 1].strip() if pi + 1 < len(parts) else ""
                rows: list[list[str]] = []
                for line in body.splitlines():
                    line = line.strip()
                    if line.startswith("|") and "|" in line[1:]:
                        cells = [c.strip() for c in line.strip("|").split("|")]
                        # Skip separator rows (e.g. |---|---|)
                        if not all(re.match(r'^[-: ]+$', c) for c in cells if c):
                            rows.append(cells)
                if rows:
                    section_data[heading] = rows
            for sheet_name in sheet_order:
                rows = section_data.get(sheet_name, [["No data available"]])
                ws = wb.create_sheet(sheet_name[:31])
                _title_rows(ws)
                ncols = max(len(r) for r in rows)
                for ri, row_data in enumerate(rows):
                    row_data = row_data + [""] * (ncols - len(row_data))
                    xl_row = ri + 5
                    style = _XL_HDR if ri == 0 else _xl_row_style(ri)
                    _xl_write_row(ws, xl_row, row_data, style)
                    ws.row_dimensions[xl_row].height = 20 if ri == 0 else 15
                for ci in range(1, ncols + 1):
                    ws.column_dimensions[get_column_letter(ci)].width = 22
                ws.freeze_panes = "A6"
        else:
            # Fallback: old 3-sheet approach when xl_map.md is absent
            ws_map = wb.create_sheet("Screen Map")
            _title_rows(ws_map)
            rows = _build_uix_screen_map(content)
            ncols = max(len(r) for r in rows)
            for ri, row_data in enumerate(rows):
                row_data = row_data + [""] * (ncols - len(row_data))
                xl_row = ri + 5
                style = _XL_HDR if ri == 0 else _xl_row_style(ri)
                _xl_write_row(ws_map, xl_row, row_data, style)
                ws_map.row_dimensions[xl_row].height = 20 if ri == 0 else 15
            for ci in range(1, ncols + 1):
                ws_map.column_dimensions[get_column_letter(ci)].width = 24
            ws_map.freeze_panes = "A6"

            ws_comp = wb.create_sheet("Angular Components")
            _title_rows(ws_comp)
            comp_rows = [
                ["Component", "Selector", "Module", "Inputs", "Outputs", "Screen(s)"],
                ["AppComponent", "app-root", "AppModule", "", "", "All"],
                ["SidenavComponent", "app-sidenav", "SharedModule", "collapsed", "navToggle", "All"],
                ["NavbarComponent", "app-navbar", "SharedModule", "title, user", "menuClick", "All"],
                ["DashboardComponent", "app-dashboard", "DashboardModule", "", "", "i.html"],
                ["KpiCardComponent", "app-kpi-card", "SharedModule", "label, value, trend", "", "i.html"],
                ["DataTableComponent", "app-data-table", "SharedModule", "rows, columns", "rowClick", "Multiple"],
                ["FormComponent", "app-form", "SharedModule", "schema, data", "formSubmit", "p7.html"],
                ["ChartComponent", "app-chart", "SharedModule", "type, data, labels", "", "p8.html, i.html"],
                ["SettingsComponent", "app-settings", "SettingsModule", "", "", "g.html"],
                ["AdminComponent", "app-admin", "AdminModule", "", "", "s.html"],
            ]
            for ri, row_data in enumerate(comp_rows):
                xl_row = ri + 5
                style = _XL_HDR if ri == 0 else _xl_row_style(ri)
                _xl_write_row(ws_comp, xl_row, row_data, style)
            for ci in range(1, 7):
                ws_comp.column_dimensions[get_column_letter(ci)].width = 22
            ws_comp.freeze_panes = "A6"

            ws_ed = wb.create_sheet("Edition Matrix")
            _title_rows(ws_ed)
            ed_rows = [
                ["Screen / Feature", "CE", "PE", "E-C", "E-O", "Notes"],
                ["Dashboard (i.html)", "Yes", "Yes", "Yes", "Yes", "Core screen"],
                ["Feature Screens (p1-p5)", "Limited", "Full", "Full", "Full", "CE has read-only"],
                ["List Screen (p6)", "Yes", "Yes", "Yes", "Yes", "Core screen"],
                ["Detail/Form (p7)", "Yes", "Yes", "Yes", "Yes", "CE limited fields"],
                ["Reports (p8)", "Basic", "Full", "Full", "Full", "CE: 3 report types"],
                ["Settings (g.html)", "Basic", "Full", "Full", "Full", "CE: system settings only"],
                ["Superadmin (s.html)", "No", "No", "Yes", "Yes", "Enterprise only"],
                ["Multi-tenancy", "No", "No", "Yes", "No", "E-C cloud only"],
                ["API Access", "No", "Yes", "Yes", "Yes", "PE and above"],
            ]
            for ri, row_data in enumerate(ed_rows):
                xl_row = ri + 5
                style = _XL_HDR if ri == 0 else _xl_row_style(ri)
                _xl_write_row(ws_ed, xl_row, row_data, style)
            for ci in range(1, 7):
                ws_ed.column_dimensions[get_column_letter(ci)].width = 20
            ws_ed.freeze_panes = "A6"

    if stage not in UIX_STAGES:
        if not tables:
            ws = wb.create_sheet("Content")
            _title_rows(ws)
            row = 5
            for t, text in tokens:
                if t in ("h1", "h2", "h3", "h4"):
                    c = ws.cell(row=row, column=1, value=_plain(text))
                    c.font = XLFont(bold=True, size=11, name="Calibri")
                elif t in ("para", "bullet", "quote") and text.strip():
                    c = ws.cell(row=row, column=1, value=_plain(text))
                    c.font = XLFont(size=10, name="Calibri")
                else:
                    row += 1
                    continue
                row += 1
            ws.column_dimensions["A"].width = 100
        else:
            table_idx = 0
            i = 0
            seen_table_starts: set[int] = set()
            while i < len(tokens) and table_idx < len(tables):
                t, line = tokens[i]
                if t == "table_row" and i not in seen_table_starts:
                    heading = f"Table {table_idx + 1}"
                    for back in range(i - 1, max(i - 8, -1), -1):
                        bt, bc = tokens[back]
                        if bt in ("h1", "h2", "h3", "h4"):
                            heading = _plain(bc)[:28]
                            break
                    sheet_name = f"{table_idx + 1}. {heading}"[:31]
                    ws = wb.create_sheet(sheet_name)
                    _title_rows(ws)
                    ws["A2"] = f"{stage_label}  —  {heading}"
                    ws["A2"].font = XLFont(size=10, color="D4A62A", name="Calibri")
                    tbl_rows = tables[table_idx]
                    ncols = max(len(r) for r in tbl_rows)
                    for ri, row_data in enumerate(tbl_rows):
                        row_data = row_data + [""] * (ncols - len(row_data))
                        xl_row = ri + 5
                        style = _XL_HDR if ri == 0 else _xl_row_style(ri)
                        _xl_write_row(ws, xl_row, row_data, style)
                        ws.row_dimensions[xl_row].height = 20 if ri == 0 else 15
                    for ci in range(1, ncols + 1):
                        col = get_column_letter(ci)
                        max_w = max(
                            len(_plain(r[ci - 1]) if ci - 1 < len(r) else "")
                            for r in tbl_rows
                        )
                        ws.column_dimensions[col].width = min(max(max_w + 2, 12), 48)
                    ws.freeze_panes = "A6"
                    seen_table_starts.add(i)
                    while i < len(tokens) and tokens[i][0] == "table_row":
                        i += 1
                    table_idx += 1
                    continue
                i += 1

    ws_info = wb.create_sheet("ProSEIT")
    ws_info["A1"] = "ProSEIT · Professional Society for Engineering, Innovation & Technology"
    ws_info["A1"].font = XLFont(bold=True, size=13, color="24364B", name="Calibri")
    for idx, row_val in enumerate([
        "Plot 5 Blue Heights Plaza, Nkrumah Road, Kampala, Uganda",
        "Tel: +256 790 334 653",
        "Website: www.proseit.org",
        "Email: info@proseit.org",
        f"Generated: {_TODAY}",
        f"Innovation: {innov_title}",
        f"Stage: {stage_label}",
    ], start=2):
        ws_info.cell(row=idx, column=1, value=row_val).font = XLFont(size=10, name="Calibri")
    ws_info.column_dimensions["A"].width = 70

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════
# PowerPoint builder
# ═══════════════════════════════════════════════════════════════════════════

def _pptx_run(tf, text: str, size: int, color: PRGBColor,
             bold: bool = False, italic: bool = False,
             align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    para = tf.add_paragraph()
    para.alignment = align
    run = para.add_run()
    run.text = _plain(text)
    run.font.size = PPt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Calibri"


def _pptx_rect(slide, left, top, width, height, fill_color: PRGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def build_pptx_doc(
    stage: str, stage_label: str, content: str, innov_title: str
) -> bytes:
    _pptx_imports()
    prs = Presentation()
    prs.slide_width = PIn(13.33)
    prs.slide_height = PIn(7.5)
    blank = prs.slide_layouts[6]

    sl = prs.slides.add_slide(blank)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = _PN
    _pptx_rect(sl, PIn(0), PIn(0), PIn(13.33), PIn(0.12), _PG)
    _pptx_rect(sl, PIn(0), PIn(7.38), PIn(13.33), PIn(0.12), _PG)

    # Logo — top-left corner of cover slide
    logo_png = _logo_png_bytes(300)
    if logo_png:
        try:
            sl.shapes.add_picture(
                io.BytesIO(logo_png),
                PIn(0.2), PIn(0.18),
                width=PIn(1.5),
            )
        except Exception:
            pass

    def _tb(left, top, width, height):
        txb = sl.shapes.add_textbox(PIn(left), PIn(top), PIn(width), PIn(height))
        txb.text_frame.word_wrap = True
        return txb.text_frame

    _pptx_run(_tb(1, 1.3, 11.33, 1.2), "PROSEIT", 54, _PG, bold=True, align=PP_ALIGN.CENTER)
    _pptx_run(_tb(1, 2.65, 11.33, 0.5),
              "Professional Society for Engineering, Innovation & Technology",
              13, _PW, align=PP_ALIGN.CENTER)
    _pptx_run(_tb(1, 3.25, 11.33, 0.3), "─" * 60, 9, _PG, align=PP_ALIGN.CENTER)
    _pptx_run(_tb(1, 3.65, 11.33, 1.0), innov_title, 26, _PW, bold=True, align=PP_ALIGN.CENTER)
    _pptx_run(_tb(1, 4.8, 11.33, 0.55), stage_label, 16, _PG, align=PP_ALIGN.CENTER)
    _pptx_run(_tb(1, 5.5, 11.33, 0.4), _TODAY, 11, _PW, align=PP_ALIGN.CENTER)
    _pptx_run(_tb(0.5, 6.75, 12.33, 0.45), PROSEIT_CONTACT, 8, _PGR, align=PP_ALIGN.CENTER)

    tokens = _parse_md(content)
    cur_heading = ""
    cur_bullets: list[str] = []
    cur_paras: list[str] = []

    def _flush(heading: str, bullets: list[str], paras: list[str]) -> None:
        if not heading and not bullets and not paras:
            return
        s = prs.slides.add_slide(blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = _PW
        _pptx_rect(s, PIn(0), PIn(0), PIn(13.33), PIn(1.2), _PN)
        _pptx_rect(s, PIn(0), PIn(1.2), PIn(13.33), PIn(0.04), _PG)
        _pptx_rect(s, PIn(0), PIn(7.15), PIn(13.33), PIn(0.04), _PG)

        txh = s.shapes.add_textbox(PIn(0.4), PIn(0.18), PIn(12.5), PIn(0.85))
        txh.text_frame.word_wrap = True
        _pptx_run(txh.text_frame, heading or stage_label, 20, _PW, bold=True)

        all_items = [("p", t) for t in paras] + [("b", t) for t in bullets]
        if not all_items:
            return

        txc = s.shapes.add_textbox(PIn(0.5), PIn(1.4), PIn(12.3), PIn(5.55))
        tfc = txc.text_frame
        tfc.word_wrap = True
        fsize = 12 if len(all_items) <= 5 else 10 if len(all_items) <= 9 else 9

        for kind, item in all_items[:14]:
            para = tfc.add_paragraph()
            run = para.add_run()
            run.text = ("•  " if kind == "b" else "") + _plain(item)
            run.font.size = PPt(fsize)
            run.font.color.rgb = _PDT
            run.font.name = "Calibri"
            para.space_before = PPt(5)

        txf = s.shapes.add_textbox(PIn(0.3), PIn(7.18), PIn(12.7), PIn(0.28))
        _pptx_run(txf.text_frame, PROSEIT_CONTACT, 7, _PGR)

    for t, cnt in tokens:
        if t in ("h1", "h2", "h3"):
            _flush(cur_heading, cur_bullets, cur_paras)
            cur_heading = _plain(cnt)
            cur_bullets = []
            cur_paras = []
        elif t == "bullet":
            cur_bullets.append(cnt)
        elif t in ("para", "quote", "h4") and cnt.strip():
            cur_paras.append(cnt)

    _flush(cur_heading, cur_bullets, cur_paras)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════════════════════════
# ZIP builders
# ═══════════════════════════════════════════════════════════════════════════

def build_stage_zip(
    stage: str, stage_label: str, content: str, innov_title: str
) -> bytes:
    slug = stage_label.replace(" · ", "_").replace(" ", "_").replace("/", "-")
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(f"{slug}.docx", build_word_doc(stage, stage_label, content, innov_title))
        if stage in EXCEL_STAGES:
            zf.writestr(f"{slug}.xlsx", build_excel_doc(stage, stage_label, content, innov_title))
        if stage in PPTX_STAGES:
            zf.writestr(f"{slug}.pptx", build_pptx_doc(stage, stage_label, content, innov_title))
        if stage in UIX_STAGES:
            uix_buf = build_uix_zip(stage, stage_label, content, innov_title)
            with zipfile.ZipFile(io.BytesIO(uix_buf), "r") as inner:
                for name in inner.namelist():
                    zf.writestr(name, inner.read(name))
    return buf.getvalue()


def build_full_zip(
    results: dict[str, str],
    stage_labels: dict[str, str],
    innov_title: str,
) -> bytes:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for stage, content in results.items():
            label = stage_labels.get(stage, stage)
            slug = label.replace(" · ", "_").replace(" ", "_").replace("/", "-")
            folder = f"{slug}/"
            zf.writestr(folder + f"{slug}.docx",
                        build_word_doc(stage, label, content, innov_title))
            if stage in EXCEL_STAGES:
                zf.writestr(folder + f"{slug}.xlsx",
                            build_excel_doc(stage, label, content, innov_title))
            if stage in PPTX_STAGES:
                zf.writestr(folder + f"{slug}.pptx",
                            build_pptx_doc(stage, label, content, innov_title))
            if stage in UIX_STAGES:
                uix_buf = build_uix_zip(stage, label, content, innov_title)
                with zipfile.ZipFile(io.BytesIO(uix_buf), "r") as inner:
                    for name in inner.namelist():
                        zf.writestr(folder + name, inner.read(name))
        index = [
            f"# {innov_title} — ProSEIT Innovation Package",
            f"Generated: {_TODAY}",
            "",
            f"Contact: {PROSEIT_CONTACT}",
            "",
            "## Contents",
        ]
        for stage, label in stage_labels.items():
            if stage in results:
                deliverables = STAGE_DELIVERABLES.get(stage, ["Word (.docx)"])
                index.append(f"- **{label}** — {', '.join(deliverables)}")
        zf.writestr("INDEX.md", "\n".join(index))
    return buf.getvalue()
